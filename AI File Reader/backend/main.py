import uuid
import traceback
from fastapi import FastAPI, File, UploadFile, HTTPException, Depends, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from fastapi.encoders import jsonable_encoder
from pydantic import BaseModel
from middlewares.auth import auth
import os
import io
import PyPDF2
import logging
from datetime import datetime
from dotenv import load_dotenv
from configs.database import get_db_connection, init_db
import openai
import sys
import json
from langchain_openai import OpenAIEmbeddings
from langchain_qdrant import QdrantVectorStore
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
import re
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
import xlrd
import requests

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

load_dotenv()

if not os.getenv("OPENAI_API_KEY"):
    print("Warning: OPENAI_API_KEY is not set. Qdrant embeddings will fail.")

app = FastAPI(title="AI File Reader API", version="1.0.0")

origins = [
    "https://questrion.vercel.app",
    "https://questrionai.ddnsfree.com"
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

init_db()

gemini_api_key = os.getenv("GEMINI_API_KEY")
if not gemini_api_key:
    print("Warning: GEMINI_API_KEY environment variable is not set")

AI = openai.OpenAI(
    api_key=gemini_api_key,
    base_url="https://generativelanguage.googleapis.com/v1beta/openai/",
)

class ChatRequest(BaseModel):
    message: str
    conversationId: str

class Creation(BaseModel):
    id: str = None
    user_id: str
    prompt: str
    content: str
    type: str
    pdf_content: str = None
    created_at: str = None


def clean_text(text: str) -> str:
    """
    Repair common PDF-extraction spacing artifacts while preserving real spaces.

    Heuristics:
    - Join spaces that appear *inside* words (lowercase-letter SPACE lowercase-letter),
      but DO NOT collapse 'a lot' or 'A lot' (exclude word-boundary + 'a' cases).
    - Keep 'New York' (right side uppercase not merged).
    - Normalize spaces around hyphens and punctuation.
    - Collapse excessive spaces while preserving newlines.
    """
    if not text:
        return text

    text = text.replace("\r\n", "\n").replace("\r", "\n")

    text = re.sub(r'(?<=\w)-\s*\n\s*(?=\w)', '-', text)

    text = re.sub(r'(?<!\b[aA])(?<=[a-z])\s+(?=[a-z])', '', text)

    text = re.sub(r'(?<=\w)\s*-\s*(?=\w)', '-', text)

    text = re.sub(r'([.,;:!?])(?=\S)', r'\1 ', text)

    text = re.sub(r'[ \t]{2,}', ' ', text)

    return text


def extract_text_from_file(file_content: bytes, filename: str) -> str:
    """Extract text from various file formats"""
    file_ext = filename.lower().split('.')[-1]
    
    if file_ext == 'pdf':
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() or ""
        return text
    
    elif file_ext in ['docx']:
        doc = DocxDocument(io.BytesIO(file_content))
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    
    elif file_ext in ['pptx']:
        prs = Presentation(io.BytesIO(file_content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
    elif file_ext in ['xlsx']:
        workbook = openpyxl.load_workbook(io.BytesIO(file_content))
        text = ""
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    text += row_text + "\n"
        return text
    
    elif file_ext in ['xls']:
        workbook = xlrd.open_workbook(file_contents=file_content)
        text = ""
        for sheet_idx in range(workbook.nsheets):
            sheet = workbook.sheet_by_index(sheet_idx)
            for row_idx in range(sheet.nrows):
                row_text = " ".join([str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)])
                if row_text.strip():
                    text += row_text + "\n"
        return text
    
    elif file_ext in ['txt', 'md']:
        return file_content.decode('utf-8', errors='ignore')
    
    else:
        raise ValueError(f"Unsupported file format: {file_ext}")

@app.get("/")
async def root():
    return {"message": "Server is Live!"}

@app.post("/api/ai/upload")
async def upload_file(
    request: Request,
    file: UploadFile = File(...),
    user_id: str = Depends(auth)
):
    try:
        logger.info(f"Processing file upload for user: {user_id}")

        allowed_extensions = ['pdf', 'docx', 'pptx', 'xlsx', 'xls', 'txt', 'md']
        file_ext = file.filename.lower().split('.')[-1]
        if file_ext not in allowed_extensions:
            raise HTTPException(
                status_code=400, 
                detail=f"File type not supported. Allowed types: {', '.join(allowed_extensions)}"
            )
        
        content = await file.read()
        await file.seek(0)
        if len(content) > 10 * 1024 * 1024:
            raise HTTPException(status_code=400, detail="File size exceeds allowed size (10MB)")
        
        try:
            file_text = extract_text_from_file(content, file.filename)
            if not file_text.strip():
                raise HTTPException(status_code=400, detail="Could not extract text from file.")
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing file: {str(e)}")

        safe_text = clean_text(file_text.replace('\x00', ''))
        
        new_id = str(uuid.uuid4())

        collection_name = f"file_{new_id}"
        try:
            if not os.getenv("OPENAI_API_KEY"):
                raise HTTPException(status_code=500, detail="OPENAI_API_KEY is required for embeddings")
            
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=400)
            docs = [Document(page_content=safe_text)]
            split_docs = text_splitter.split_documents(docs)
            
            if not split_docs or not any(doc.page_content.strip() for doc in split_docs):
                raise HTTPException(status_code=400, detail="No meaningful content found in file for indexing")
            
            embeddings = OpenAIEmbeddings(model="text-embedding-3-small")
            QdrantVectorStore.from_documents(
                documents=split_docs,
                url="http://vector-db:6333",
                collection_name=collection_name,
                embedding=embeddings
            )
            logger.info(f"Successfully indexed {len(split_docs)} chunks into Qdrant collection {collection_name}")
        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"Vector indexing failed: {e}")
            raise HTTPException(status_code=500, detail=f"Failed to create vector index: {str(e)}")

        conn = get_db_connection()
        cur = conn.cursor()
        try:
            logger.info(f"Inserting creation with id={new_id}, user_id={user_id}")
            cur.execute("""
                INSERT INTO creations (id, user_id, prompt, content, type, pdf_content, created_at)
                VALUES (%s, %s, %s, %s, %s, %s, NOW())
                RETURNING id
            """, (
                new_id,
                user_id,
                file.filename,
                '',
                'file',
                safe_text
            ))
            creation_id = cur.fetchone()['id']
            conn.commit()
        except Exception as e:
            logger.error(f"DB insert error: {e}")
            logger.error(traceback.format_exc())
            try:
                delete_url = f"http://vector-db:6333/collections/{collection_name}"
                requests.delete(delete_url)
            except Exception:
                pass
            raise HTTPException(status_code=500, detail="Failed to save file information")
        finally:
            cur.close()
            conn.close()

        return JSONResponse(content={
            "success": True,
            "conversationId": new_id,
            "fileText": safe_text,
            "extractedText": safe_text,
            "pdfText": safe_text       
        })
    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"Error in /api/ai/upload: {e}")
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/ai/pdf")
async def upload_pdf(
    request: Request,
    pdf: UploadFile = File(...),
    user_id: str = Depends(auth)
):
    return await upload_file(request, pdf, user_id)

@app.post("/api/ai/chat")
async def chat_with_file(request: ChatRequest, user_id: str = Depends(auth)):
    try:
        logger.info(f"Processing chat request for user: {user_id}, conversation: {request.conversationId}")
        
        conn = get_db_connection()
        cur = conn.cursor()
        try:
            cur.execute("""
                SELECT pdf_content FROM creations 
                WHERE user_id = %s AND id = %s
            """, (user_id, request.conversationId))
            creation = cur.fetchone()
            if not creation:
                raise HTTPException(status_code=404, detail="Conversation not found")

            file_content = creation['pdf_content'] or ''
            file_content = clean_text(file_content)
        finally:
            cur.close()
            conn.close()

        try:
            if not os.getenv("OPENAI_API_KEY"):
                raise Exception("OPENAI_API_KEY is required for embeddings")
            embeddings = OpenAIEmbeddings(model="text-embedding-3-large")
            vector_store = None
            last_error = None
            for prefix in ["file_", "pdf_"]:
                collection_name = f"{prefix}{request.conversationId}"
                try:
                    vector_store = QdrantVectorStore.from_existing_collection(
                        url="http://vector-db:6333",
                        collection_name=collection_name,
                        embedding=embeddings
                    )
                    break
                except Exception as e:
                    last_error = e
            
            if vector_store is None:
                raise HTTPException(
                    status_code=409,
                    detail="Vector index missing or unreachable. Please re-upload the file."
                )
            
            matches = vector_store.similarity_search(query=request.message, k=5)
            if not matches:
                raise HTTPException(
                    status_code=409,
                    detail="No vector context found for this conversation. Please re-upload the file."
                )
            retrieved_context = "\n\n".join(doc.page_content for doc in matches)
            retrieved_context = clean_text(retrieved_context)
        except HTTPException:
            raise
        except Exception:
            raise HTTPException(
                status_code=409,
                detail="Vector index missing or unreachable. Please re-upload the file."
            )

        SYSTEM_PROMPT = f"""
            You are a helpful AI Assistant who answers user queries based only on the retrieved context.
            
            Context:
            {retrieved_context}
            
            User Question: {request.message}
            
            Answer concisely and cite only from the context and give the page number of the answer if the file contains page number.
        """

        try:
            response = AI.chat.completions.create(
                model="gemini-2.0-flash",
                messages=[{"role": "user", "content": SYSTEM_PROMPT}],
                temperature=0.7,
                max_tokens=1000,
            )
            if not response.choices or not hasattr(response.choices[0], 'message'):
                raise Exception("Invalid response from AI API")
            ai_content = response.choices[0].message.content
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Error calling AI API: {str(e)}")

        conn = get_db_connection()
        cur = conn.cursor()
        try:
            cur.execute("""
                INSERT INTO chat_messages (user_id, conversation_id, message, response)
                VALUES (%s, %s, %s, %s)
            """, (user_id, request.conversationId, request.message, ai_content))
            conn.commit()
        finally:
            cur.close()
            conn.close()

        return JSONResponse(content={
            "success": True,
            "content": ai_content
        })

    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"Error in /api/ai/chat: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/user/get-user-creations")
async def get_user_creations(
    request: Request,
    user_id: str = Depends(auth)
):
    try:
        conn = get_db_connection()
        cur = conn.cursor()
        try:
            cur.execute("""
                SELECT id, prompt, content, type, created_at, publish
                FROM creations
                WHERE user_id = %s
                ORDER BY created_at DESC
            """, (user_id,))
            creations = cur.fetchall()

            creations_with_chats = []
            for creation in creations:
                cur.execute("""
                    SELECT message, response, created_at
                    FROM chat_messages
                    WHERE conversation_id = %s AND user_id = %s
                    ORDER BY created_at ASC
                """, (creation['id'], user_id))
                chat_messages = cur.fetchall()

                creation_dict = dict(creation)
                creation_dict['chat_messages'] = chat_messages
                creations_with_chats.append(creation_dict)
        finally:
            cur.close()
            conn.close()

        return JSONResponse(content=jsonable_encoder({
            "success": True,
            "creations": creations_with_chats
        }))
    except Exception as e:
        logger.error(f"Error fetching user creations: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/user/delete-creation/{creation_id}")
async def delete_creation(
    creation_id: str,
    user_id: str = Depends(auth)
):
    try:
        logger.info(f"Deleting creation {creation_id} for user {user_id}")
        
        conn = get_db_connection()
        cur = conn.cursor()
        try:
            cur.execute("""
                SELECT id FROM creations 
                WHERE id = %s AND user_id = %s
            """, (creation_id, user_id))
            
            if not cur.fetchone():
                raise HTTPException(status_code=404, detail="Creation not found")

            cur.execute("""
                DELETE FROM creations 
                WHERE id = %s AND user_id = %s
            """, (creation_id, user_id))
            
            conn.commit()
            
            try:
                qdrant_url = "http://vector-db:6333"
                for prefix in ["file_", "pdf_"]:
                    collection_name = f"{prefix}{creation_id}"
                    delete_url = f"{qdrant_url}/collections/{collection_name}"
                    response = requests.delete(delete_url)
                    if response.status_code == 200:
                        logger.info(f"Successfully deleted Qdrant collection: {collection_name}")
                    elif response.status_code == 404:
                        logger.info(f"Qdrant collection {collection_name} not found (already deleted)")
                    else:
                        logger.warning(f"Failed to delete Qdrant collection {collection_name}: {response.status_code}")
            except Exception as e:
                logger.warning(f"Error deleting Qdrant collection(s): {e}")
                
        finally:
            cur.close()
            conn.close()

        return JSONResponse(content={
            "success": True,
            "message": "Creation deleted successfully"
        })
        
    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"Error deleting creation: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/ai/suggestions")
async def get_suggestions(conversationId: str, user_id: str = Depends(auth)):
    try:
        conn = get_db_connection()
        cur = conn.cursor()
        try:
            cur.execute("""
                SELECT pdf_content FROM creations
                WHERE user_id = %s AND id = %s
            """, (user_id, conversationId))
            creation = cur.fetchone()
            if not creation:
                raise HTTPException(status_code=404, detail="Conversation not found")
            file_content = creation['pdf_content'] or ''
            file_content = clean_text(file_content)

            cur.execute("""
                SELECT message FROM chat_messages
                WHERE user_id = %s AND conversation_id = %s AND message IS NOT NULL
                ORDER BY created_at ASC
            """, (user_id, conversationId))
            asked_rows = cur.fetchall() or []
        finally:
            cur.close()
            conn.close()

        asked_questions = [r['message'] for r in asked_rows]

        context = clean_text((file_content or '')[:4000])
        prev_qs_text = "\n".join(f"- {q}" for q in asked_questions[:20])

        system_prompt = "Return only a JSON array of strings with 5 concise, specific suggested questions about the document."
        user_prompt = f"""Generate 5 short, specific questions a user could ask about the document below and which answers are available in the file.
                        Do NOT repeat or paraphrase any of these already asked questions:
                        {prev_qs_text or "(none)"}

                        Return ONLY a JSON array of strings and nothing else.

                        Document:
                        {context}
                    """

        try:
            response = AI.chat.completions.create(
                model="gemini-2.0-flash",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt},
                ],
                temperature=0.7,
                max_tokens=300,
            )
            raw = response.choices[0].message.content.strip()

            if raw.startswith("```"):
                raw = re.sub(r"^```(?:json)?\s*|\s*```$", "", raw, flags=re.DOTALL).strip()

            suggestions = json.loads(raw)
            if not isinstance(suggestions, list):
                raise ValueError("Model did not return a list")

            def normalize(s: str) -> str:
                return re.sub(r"\s+", " ", re.sub(r"[^\w\s]", "", (s or "").lower())).strip()

            asked_set = {normalize(q) for q in asked_questions}
            suggestions = [str(x) for x in suggestions if normalize(str(x)) not in asked_set][:6]
        except Exception:
            suggestions = [
                "What is the main purpose of this document?",
                "Can you summarize the key points?",
                "Which section covers the most important topic?",
                "What steps or processes are described?",
                "Are there any key terms or definitions I should know?"
            ]

        return JSONResponse(content={"success": True, "suggestions": suggestions})
    except HTTPException as he:
        raise he
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=5000)
